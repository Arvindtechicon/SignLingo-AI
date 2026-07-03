import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Page setup - 16:9 widescreen layout
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Styling colors
    INDIGO = RGBColor(79, 70, 229)    # Theme primary
    SLATE = RGBColor(15, 23, 42)      # Deep text
    MUTED = RGBColor(100, 116, 139)   # Subtitle / grey
    EMERALD = RGBColor(16, 185, 129)  # Accent green
    
    # ----------------- SLIDE 1: Title Slide -----------------
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only layout
    
    # Main Title
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SignLingo-AI: Milestone 2"
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = INDIGO
    p.font.name = 'Outfit'
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "Real-Time Gesture Recognition & Sign Language Assessment"
    p2.font.size = Pt(22)
    p2.font.color.rgb = SLATE
    p2.font.name = 'Inter'
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(12)
    
    # Presenter Details
    p3 = tf.add_paragraph()
    p3.text = "Presented by: Arvindtechicon | Project Milestone 2 Deliverables"
    p3.font.size = Pt(14)
    p3.font.color.rgb = MUTED
    p3.font.name = 'Inter'
    p3.space_before = Pt(40)
    
    # ----------------- Helper for Content Slides -----------------
    def add_content_slide(title, bullet_points):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Slide Title
        tBox = s.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
        tf_t = tBox.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(32)
        p_t.font.bold = True
        p_t.font.color.rgb = INDIGO
        p_t.font.name = 'Outfit'
        
        # Slide Content Box
        cBox = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.8))
        tf_c = cBox.text_frame
        tf_c.word_wrap = True
        
        first = True
        for bp in bullet_points:
            if first:
                p_c = tf_c.paragraphs[0]
                first = False
            else:
                p_c = tf_c.add_paragraph()
            
            p_c.text = bp[0]
            p_c.font.size = Pt(bp[1])
            p_c.font.name = 'Inter'
            p_c.font.color.rgb = SLATE
            p_c.space_before = Pt(bp[2])
            
            # Highlight keyword bolding if colon is present
            if ":" in bp[0]:
                parts = bp[0].split(":", 1)
                p_c.text = ""
                run_bold = p_c.add_run()
                run_bold.text = parts[0] + ":"
                run_bold.font.bold = True
                run_bold.font.color.rgb = SLATE
                
                run_norm = p_c.add_run()
                run_norm.text = parts[1]
                run_norm.font.bold = False
                run_norm.font.color.rgb = SLATE
                
            if bp[3] > 0:
                p_c.level = bp[3]
                p_c.font.color.rgb = MUTED
        
        return s

    # ----------------- SLIDE 2: Objectives & Scope -----------------
    add_content_slide(
        "Objectives & Scope of Milestone 2",
        [
            ("Develop a robust hand-tracking frontend and pattern recognition assessment backend.", 18, 12, 0),
            ("Real-Time Tracking: Integrate MediaPipe client-side to track multi-hand and upper body posture.", 16, 10, 1),
            ("Dynamic Assessment: Build algorithms capable of measuring coordinate velocity and sequence alignment.", 16, 6, 1),
            ("Storage Ingestion: Capture practice history alongside raw coordinate trails directly into user records.", 16, 6, 1),
            ("User-Friendly HUD: Surface accuracy ratings and corrective finger joint instructions visually in real-time.", 16, 6, 1)
        ]
    )

    # ----------------- SLIDE 3: System Architecture & Data Pipeline -----------------
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Title
    tBox = s3.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
    p_t = tBox.text_frame.paragraphs[0]
    p_t.text = "System Architecture & Data Pipeline"
    p_t.font.size = Pt(32)
    p_t.font.bold = True
    p_t.font.color.rgb = INDIGO
    p_t.font.name = 'Outfit'
    
    # Description
    dBox = s3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.0))
    p_d = dBox.text_frame.paragraphs[0]
    p_d.text = "The platform features a decoupled real-time pipeline that processes video frames client-side and runs similarity metrics in the backend:"
    p_d.font.size = Pt(16)
    p_d.font.name = 'Inter'
    p_d.font.color.rgb = SLATE
    
    # Draw Flowchart Boxes
    steps = [
        "Webcam Input\n(Video Frame)",
        "MediaPipe WASM\n(Joint Tracking)",
        "WebSockets\n(JSON Ingestion)",
        "FastAPI Deque\n(Frame Buffering)",
        "Assessment\n(Cosine/Procrustes/DTW)",
        "SQLite Storage\n(Attempt Logs)"
    ]
    
    box_width = Inches(1.6)
    box_height = Inches(1.1)
    top_pos = Inches(3.4)
    left_margin = Inches(0.8)
    spacing = Inches(0.4)
    
    for i, step in enumerate(steps):
        left_pos = left_margin + i * (box_width + spacing)
        # Add box shape
        shape = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = INDIGO if i % 2 == 0 else SLATE
        shape.line.color.rgb = INDIGO
        shape.line.width = Pt(1.5)
        
        tf_s = shape.text_frame
        tf_s.word_wrap = True
        p_s = tf_s.paragraphs[0]
        p_s.text = step
        p_s.font.size = Pt(11)
        p_s.font.name = 'Inter'
        p_s.font.bold = True
        p_s.font.color.rgb = RGBColor(255, 255, 255)
        p_s.alignment = PP_ALIGN.CENTER
        
        # Add arrow connector (if not the last box)
        if i < len(steps) - 1:
            arrow_left = left_pos + box_width + Inches(0.05)
            arrow_width = spacing - Inches(0.1)
            arrow_top = top_pos + box_height / 2 - Inches(0.1)
            arrow = s3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, arrow_width, Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.fill.background()

    # ----------------- SLIDE 4: Frontend Implementation -----------------
    add_content_slide(
        "Frontend Tracking & UI Features",
        [
            ("Key UI features built inside the Practice Sandbox interface:", 18, 12, 0),
            ("Custom Canvas Overlay: Draws Cobalt Blue hand skeletons and Emerald Green joints over the camera feed.", 16, 10, 1),
            ("Visual Pose Overlay: Renders Violet indicators representing Shoulders, Elbows, and Arm movements.", 16, 6, 1),
            ("Coordinate Mirroring: Flips the horizontal draw axis to match the mirrored camera display naturally.", 16, 6, 1),
            ("On-Demand Camera Toggle: Built an Enable/Disable button with clean video camera placeholders.", 16, 6, 1),
            ("Sign Selector: Full alphabet dropdown (A-Z) and 13 dynamic gesture sequence selections.", 16, 6, 1)
        ]
    )

    # ----------------- SLIDE 5: Pattern Recognition Algorithms -----------------
    add_content_slide(
        "AI Pattern Recognition & Assessment Models",
        [
            ("Integrated three core evaluation algorithms to determine gesture correctness:", 18, 12, 0),
            ("Cosine Similarity: Measures the vector angular direction correlation for static letters.", 16, 10, 1),
            ("Procrustes Shape Comparison: Computes coordinate disparity, normalizing translation, scaling, and rotations.", 16, 6, 1),
            ("Dynamic Time Warping (DTW): Aligns motion sequences of variable speed/length for dynamic words.", 16, 6, 1),
            ("Corrective Vector Checking: Calculates pointing angles (MCP to Tip) to return detailed warnings (e.g. thumb/pinky misalignment).", 16, 6, 1)
        ]
    )

    # ----------------- SLIDE 6: Relational Database Schema -----------------
    add_content_slide(
        "Database Architecture & SQLite Fallbacks",
        [
            ("Robust storage configuration mapping user performance directly in user profiles:", 18, 12, 0),
            ("AttemptLog Table: Stores scores, timestamps, duration, correctness, and corrective feedbacks.", 16, 10, 1),
            ("JSONB Coordinate Trails: Saves the entire coordinate history of the practice attempt for review.", 16, 6, 1),
            ("SQLite Connection Fallback: Automatically detects if PostgreSQL is offline and switches to local test.db.", 16, 6, 1),
            ("Auto-Migrations: Runs database metadata create_all scripts dynamically to instantiate schemas.", 16, 6, 1)
        ]
    )

    # ----------------- SLIDE 7: Verification & Testing Pipeline -----------------
    add_content_slide(
        "Testing, Simulation & Verification Pipelines",
        [
            ("Three-layer verification pipeline ensuring system stability under all scenarios:", 18, 12, 0),
            ("Automated Pytest Suite: Runs 5 backend tests checking OAuth registration, token login, and WebSocket loops.", 16, 10, 1),
            ("Template Validator: Checks memory loading and parsing of reference coordinate sets.", 16, 6, 1),
            ("DTW Stream Simulation: Simulates streaming perfect templates frame-by-frame to verify 100% evaluation scores.", 16, 6, 1),
            ("Compilation Assurance: Successful production bundling via Vite npm build.", 16, 6, 1)
        ]
    )

    # ----------------- SLIDE 8: Technical Challenges Resolved -----------------
    add_content_slide(
        "Engineering Challenges Faced & Overcome",
        [
            ("Key technical hurdles resolved during development:", 18, 12, 0),
            ("WebAssembly Clashing: Solved by sequential execution of models, preventing Emscripten thread conflicts.", 16, 10, 1),
            ("Overlay Misalignment: Corrected by adding coordinate mirroring formulas to match CSS camera scale flips.", 16, 6, 1),
            ("NumPy JSON Serialization: Cast NumPy float64 and bool types to native Python structures before WebSocket streaming.", 16, 6, 1),
            ("Procrustes Singular Exceptions: Handled degenerate matrices safely to keep the backend online.", 16, 6, 1)
        ]
    )

    # ----------------- SLIDE 9: Outcomes & Summary -----------------
    add_content_slide(
        "Outcomes & Deliverables Summary",
        [
            ("Milestone 2 is complete, verified, and successfully pushed to main branch:", 18, 12, 0),
            ("Interactive Sandbox: Real-time hand-body overlays, accuracy wheels, and on/off camera buttons.", 16, 10, 1),
            ("Assessment Engine: Complete Cosine, Procrustes, and DTW matching logic.", 16, 6, 1),
            ("Untracked Local Configs: Excluded system documents (e.g. presentations, PDF reports) and local .env keys.", 16, 6, 1),
            ("100% Passing Tests: Full coverage on auth gateways and WebSocket connections.", 16, 6, 1)
        ]
    )

    # ----------------- SLIDE 10: Milestone 3 Roadmap -----------------
    add_content_slide(
        "Milestone 3 Roadmap: Next Steps",
        [
            ("Looking Ahead: I will implement and execute all Milestone 3 features to deliver a complete pipeline.", 18, 12, 0),
            ("Interactive Practice Modules: Integrate gamified lessons and custom progress indicators.", 16, 10, 1),
            ("Advanced Analytics Dashboards: Build historical stats charts, time-spent counters, and progress reports.", 16, 6, 1),
            ("Assessment Enhancements: Expand the dictionary template libraries and refine accuracy matching thresholds.", 16, 6, 1),
        ]
    )

    # Save to workspace root
    output_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..", "milestone_2_presentation_v2.pptx"))
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")


if __name__ == "__main__":
    create_presentation()
